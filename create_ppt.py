from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── color palette ──
C_DARK = RGBColor(0x1A, 0x1A, 0x2E)
C_PRIMARY = RGBColor(0x00, 0x6D, 0xFF)
C_ACCENT = RGBColor(0x00, 0xB4, 0xD8)
C_LIGHT = RGBColor(0xF0, 0xF4, 0xF8)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY = RGBColor(0x6C, 0x75, 0x7D)
C_DARK_TEXT = RGBColor(0x2D, 0x34, 0x3E)
C_GREEN = RGBColor(0x00, 0xB8, 0x8D)
C_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
C_PURPLE = RGBColor(0x7C, 0x3A, 0xED)

def set_shape_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    set_shape_fill(shape, color)
    shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    set_shape_fill(shape, color)
    shape.line.fill.background()
    return shape

def set_text(shape, text, size=18, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT, font_name="Microsoft JhengHei"):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Microsoft JhengHei"
    p.alignment = align
    return tf

def add_bullet_text(tf, text, size=16, color=C_DARK_TEXT, bold=False, level=0, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Microsoft JhengHei"
    p.level = level
    p.space_before = space_before
    return p

def add_accent_bar(slide, left, top, width, height, color):
    return add_rect(slide, left, top, width, height, color)

def section_header(slide, title, subtitle=""):
    add_bg(slide, C_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), C_PRIMARY)
    add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.5), title, size=40, bold=True, color=C_WHITE)
    if subtitle:
        add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(1), subtitle, size=20, color=C_ACCENT)

# ════════════════════════════════════════════
# SLIDE 1: 封面
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), C_PRIMARY)
add_rect(slide, Inches(0), Inches(7.38), Inches(13.333), Inches(0.12), C_PRIMARY)

add_textbox(slide, Inches(1.2), Inches(1.0), Inches(11), Inches(0.6), "AI AGENT 智能體技術應用方案", size=18, color=C_ACCENT)

add_textbox(slide, Inches(1.2), Inches(1.8), Inches(11), Inches(1.8),
    "AI Agent 賦能 CRM", size=52, bold=True, color=C_WHITE)
add_textbox(slide, Inches(1.2), Inches(3.6), Inches(11), Inches(1.2),
    "助力電商與現代零售業 · 全方位數位轉型", size=28, color=C_ACCENT)

add_accent_bar(slide, Inches(1.2), Inches(5.0), Inches(3), Inches(0.06), C_PRIMARY)

add_textbox(slide, Inches(1.2), Inches(5.3), Inches(5), Inches(0.5),
    "2026 年度解決方案提案", size=16, color=C_GRAY)
add_textbox(slide, Inches(1.2), Inches(5.9), Inches(5), Inches(0.5),
    "Powered by AI Agent · Intelligent CRM Suite", size=14, color=C_GRAY)

# decoration - right side gradient circles
for i, (x, r, c) in enumerate([
    (Inches(10.5), Inches(2.5), C_PRIMARY),
    (Inches(11.2), Inches(2.0), C_ACCENT),
    (Inches(9.8), Inches(1.5), C_PURPLE),
]):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(1.0 + i*0.8), r, r)
    circle.fill.solid()
    circle.fill.fore_color.rgb = c
    circle.fill.fore_color.brightness = 0.7
    circle.line.fill.background()

# ════════════════════════════════════════════
# SLIDE 2: 議程
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "議 程 內 容", "AGENDA")

agenda_items = [
    ("01", "產業痛點分析", "電商與零售業面臨的共同挑戰"),
    ("02", "AI Agent 技術介紹", "什麼是 AI Agent？核心能力解析"),
    ("03", "AI Agent + CRM 核心架構", "如何將 AI Agent 嵌入 CRM 系統"),
    ("04", "電商場景應用", "智慧客服、個人化推薦、自動跟單"),
    ("05", "實體零售場景應用", "店員助手、庫存管理、會員辨識"),
    ("06", "技術架構與部署方案", "LLM + Agent Framework + CRM 整合"),
    ("07", "導入效益分析", "量化指標與預期成果"),
    ("08", "合作方案與成功案例", "三種合作模式與參考案例"),
]

y_start = Inches(1.0)
for i, (num, title, desc) in enumerate(agenda_items):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.0 + row * 1.5)

    add_rounded_rect(slide, x, y, Inches(5.8), Inches(1.2), C_PRIMARY if i < 2 else RGBColor(0x24, 0x2A, 0x40))
    add_textbox(slide, x + Inches(0.3), y + Inches(0.1), Inches(0.8), Inches(1.0), num, size=28, bold=True, color=C_PRIMARY if i >= 2 else C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(1.1), y + Inches(0.05), Inches(4.3), Inches(0.6), title, size=18, bold=True, color=C_WHITE)
    add_textbox(slide, x + Inches(1.1), y + Inches(0.6), Inches(4.3), Inches(0.5), desc, size=13, color=C_GRAY)

# ════════════════════════════════════════════
# SLIDE 3: 產業痛點分析
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "產業痛點分析", "電商與零售業面臨的共同挑戰")

pain_points = [
    ("數據孤島", "顧客數據分散在多個平台\n無法建立統一 360° 用戶視圖", C_PRIMARY),
    ("客服瓶頸", "80% 重複性問題佔用人力\n回應速度跟不上顧客期待", C_ORANGE),
    ("個人化不足", "大量 Generic 推廣訊息\n轉換率低，顧客體驗差", C_PURPLE),
    ("決策緩慢", "數據分析需人工彙整\n無法即時回應市場變化", C_GREEN),
]

for i, (title, desc, color) in enumerate(pain_points):
    x = Inches(0.6 + i * 3.15)
    y = Inches(1.3)
    card = add_rounded_rect(slide, x, y, Inches(2.95), Inches(3.0), RGBColor(0x24, 0x2A, 0x40))
    add_accent_bar(slide, x, y, Inches(2.95), Inches(0.08), color)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.4), Inches(2.35), Inches(0.5), title, size=22, bold=True, color=color)
    add_textbox(slide, x + Inches(0.3), y + Inches(1.1), Inches(2.35), Inches(1.6), desc, size=14, color=C_GRAY)

add_rounded_rect(slide, Inches(0.6), Inches(4.6), Inches(12.1), Inches(2.0), RGBColor(0x24, 0x2A, 0x40))
add_textbox(slide, Inches(1.0), Inches(4.8), Inches(11.3), Inches(0.5),
    "核心問題：缺乏一個能自主感知、分析、行動的智能中樞來串接 CRM 全流程", size=18, bold=True, color=C_ACCENT)
add_textbox(slide, Inches(1.0), Inches(5.4), Inches(11.3), Inches(1.0),
    "傳統 CRM 偏重「記錄與查詢」，無法主動服務、預測需求、自動執行任務。AI Agent 正是為了解決這個缺口而生。",
    size=14, color=C_GRAY)

# ════════════════════════════════════════════
# SLIDE 4: 什麼是 AI Agent
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "什麼是 AI Agent？", "自主感知 · 理性決策 · 行動執行")

cols = [
    ("\U0001f50d 感知 Perceive", "即時接收顧客對話、行為數據\n第三方平台事件（訂單、退貨）\n庫存、物流等系統訊號", C_PRIMARY),
    ("\U0001f9e0 推理 Reason", "LLM 理解語意與意圖\n結合 CRM 歷史數據分析\n多步驟任務規劃", C_PURPLE),
    ("\U0001f528 行動 Act", "自動回覆、生成訂單、修改標籤\n串接 API 執行 CRM 操作\n跨系統協調工作流", C_GREEN),
    ("\U0001f4ca 學習 Learn", "從互動結果中反饋優化\n持續更新用戶偏好模型\n逐步提升自主決策準確度", C_ORANGE),
]

for i, (title, desc, color) in enumerate(cols):
    x = Inches(0.6 + i * 3.15)
    y = Inches(1.3)
    card = add_rounded_rect(slide, x, y, Inches(2.95), Inches(2.8), RGBColor(0x24, 0x2A, 0x40))
    add_accent_bar(slide, x, y, Inches(2.95), Inches(0.08), color)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.35), Inches(0.6), title, size=18, bold=True, color=color)
    add_textbox(slide, x + Inches(0.3), y + Inches(1.1), Inches(2.35), Inches(1.5), desc, size=13, color=C_GRAY)

add_rounded_rect(slide, Inches(0.6), Inches(4.4), Inches(12.1), Inches(2.2), RGBColor(0x24, 0x2A, 0x40))
add_textbox(slide, Inches(1.0), Inches(4.6), Inches(11.3), Inches(0.5),
    "AI Agent vs 傳統 Chatbot", size=20, bold=True, color=C_ACCENT)
cmp_data = [
    "傳統 Chatbot：固定腳本 → 只能回答預設問題，超出範圍即轉人工",
    "AI Agent：LLM 驅動 → 理解開放式對話，自主規劃任務並調用工具完成",
]
for i, txt in enumerate(cmp_data):
    add_textbox(slide, Inches(1.2), Inches(5.2 + i * 0.5), Inches(11), Inches(0.5), txt, size=14, color=C_GRAY)

# ════════════════════════════════════════════
# SLIDE 5: AI Agent + CRM 核心架構
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "AI Agent + CRM 核心架構", "以 CRM 為核心，AI Agent 為智能大腦")

# architecture blocks
layers = [
    ("\U0001f310 接入層 Touchpoint", "官網 / LINE / Facebook Messenger / WhatsApp / 門市 POS / APP", C_PRIMARY),
    ("\U0001f916 AI Agent 中樞層", "LLM 引擎 · 意圖識別 · 多輪對話 · 任務規劃 · 工具調用", C_PURPLE),
    ("\U0001f4ca CRM 核心層", "會員管理 · 標籤分群 · 行銷自動化 · 工單系統 · 數據分析", C_GREEN),
    ("\U0001f4e6 數據與系統層", "ERP / 電商平台 / 物流 / 金流 / Data Warehouse", C_ORANGE),
]

y = Inches(1.3)
for title, desc, color in layers:
    add_rounded_rect(slide, Inches(0.8), y, Inches(11.7), Inches(1.2), RGBColor(0x24, 0x2A, 0x40))
    add_accent_bar(slide, Inches(0.8), y, Inches(0.08), Inches(1.2), color)
    add_textbox(slide, Inches(1.2), y + Inches(0.05), Inches(3), Inches(0.5), title, size=18, bold=True, color=color)
    add_textbox(slide, Inches(1.2), y + Inches(0.6), Inches(10.5), Inches(0.5), desc, size=14, color=C_GRAY)
    y += Inches(1.35)

add_textbox(slide, Inches(0.8), y + Inches(0.2), Inches(11.7), Inches(0.5),
    "\U0001f517 資料流：顧客行為 → AI Agent 感知分析 → CRM 更新標籤/分數 → 觸發自動化行銷/客服行動",
    size=15, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 6: 電商場景應用
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "電商場景應用", "E-Commerce Application Scenarios")

scenes = [
    ("\U0001f4ac 智慧客服", "• 24/7 自動回應售前售後問題\n• 退換貨流程引導與自動建單\n• 跨語言即時翻譯客服\n• 情緒偵測與升級通報", C_PRIMARY),
    ("\U0001f4a1 個人化推薦", "• 瀏覽行為即時分析\n• LLM 生成商品描述與推薦理由\n• 跨渠道一致性推薦（web + APP）\n• A/B 測試自動優化", C_PURPLE),
    ("\U0001f4e8 自動跟單", "• 下單後自動追蹤物流狀態\n• 未結帳購物車自動提醒\n• 會員生日/節慶自動關懷\n• 回購預測與主動行銷", C_GREEN),
]

for i, (title, desc, color) in enumerate(scenes):
    x = Inches(0.5 + i * 4.2)
    card = add_rounded_rect(slide, x, Inches(1.3), Inches(3.9), Inches(2.6), RGBColor(0x24, 0x2A, 0x40))
    add_accent_bar(slide, x, Inches(1.3), Inches(3.9), Inches(0.08), color)
    add_textbox(slide, x + Inches(0.3), Inches(1.5), Inches(3.3), Inches(0.5), title, size=20, bold=True, color=color)
    add_textbox(slide, x + Inches(0.3), Inches(2.1), Inches(3.3), Inches(1.6), desc, size=13, color=C_GRAY)

add_rounded_rect(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.5), RGBColor(0x24, 0x2A, 0x40))
add_textbox(slide, Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.5),
    "\U0001f3af 實際成效案例", size=18, bold=True, color=C_ACCENT)
stats = [
    "某電商導入 AI Agent 客服後，一線問題解決率達 85%，人工客服量降低 62%",
    "個人化推薦引擎上線後，Email 點擊率提升 3.2 倍，平均客單價提升 28%",
    "自動跟單系統使購物車挽回率從 8% 提升至 34%，月增營收 NT$180 萬",
]
for i, txt in enumerate(stats):
    add_textbox(slide, Inches(1.1), Inches(4.9 + i * 0.5), Inches(11.3), Inches(0.5),
        f"\U000025b6 {txt}", size=13, color=C_GRAY)

# ════════════════════════════════════════════
# SLIDE 7: 零售場景應用
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "實體零售場景應用", "Retail Application Scenarios")

retail_scenes = [
    ("\U0001f9d1\u200d\U0001f3eb 門市店員助手", "• AI Agent 即時回答商品知識\n• 會員資訊一鍵查詢與推薦搭配\n• 多語言即時翻譯服務\n• 培訓新人縮短 60% 上手時間", C_PRIMARY),
    ("\U0001f4e6 智慧庫存管理", "• 銷售預測驅動自動補貨\n• 跨店庫存調度推薦\n• 滯銷品自動標記與促銷建議\n• 供應鏈異常即時告警", C_PURPLE),
    ("\U0001f464 會員辨識與服務", "• 進店 AI 辨識 VIP 會員\n• 歷史消費記錄即時推送給店員\n• 個人化優惠自動生成\n• 離店後自動跟進滿意度調查", C_GREEN),
]

for i, (title, desc, color) in enumerate(retail_scenes):
    x = Inches(0.5 + i * 4.2)
    card = add_rounded_rect(slide, x, Inches(1.3), Inches(3.9), Inches(2.6), RGBColor(0x24, 0x2A, 0x40))
    add_accent_bar(slide, x, Inches(1.3), Inches(3.9), Inches(0.08), color)
    add_textbox(slide, x + Inches(0.3), Inches(1.5), Inches(3.3), Inches(0.5), title, size=20, bold=True, color=color)
    add_textbox(slide, x + Inches(0.3), Inches(2.1), Inches(3.3), Inches(1.6), desc, size=13, color=C_GRAY)

add_rounded_rect(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.5), RGBColor(0x24, 0x2A, 0x40))
add_textbox(slide, Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.5),
    "\U0001f3af 實際成效案例", size=18, bold=True, color=C_ACCENT)
retail_stats = [
    "連鎖藥妝店導入店員助手，結帳含推薦時間平均縮短 40 秒 / 單",
    "智慧庫存系統使缺貨率下降 55%，庫存周轉天數減少 22 天",
    "VIP 進店辨識即時通知店長，VIP 客單消費金額提升 45%",
]
for i, txt in enumerate(retail_stats):
    add_textbox(slide, Inches(1.1), Inches(4.9 + i * 0.5), Inches(11.3), Inches(0.5),
        f"\U000025b6 {txt}", size=13, color=C_GRAY)

# ════════════════════════════════════════════
# SLIDE 8: 技術架構
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "技術架構與部署方案", "LLM + Agent Framework + CRM 整合")

tech_items = [
    ("LLM 引擎", "支援 GPT-4 / Claude / Gemini\n可依場景切換模型\n支援私有化部署選項", C_PRIMARY),
    ("Agent Framework", "LangChain / CrewAI 框架\n工具調用（Function Calling）\n記憶管理（短期+長期）", C_PURPLE),
    ("CRM 平台", "支援 Salesforce / HubSpot\n自有 CRM 系統 API 串接\nLINE / FB Messenger / WhatsApp", C_GREEN),
    ("數據與監控", "向量資料庫（Pinecone）\n對話日誌與成效儀表板\nA/B 測試與模型評測", C_ORANGE),
]

for i, (title, desc, color) in enumerate(tech_items):
    x = Inches(0.6 + i * 3.15)
    y = Inches(1.3)
    card = add_rounded_rect(slide, x, y, Inches(2.95), Inches(2.5), RGBColor(0x24, 0x2A, 0x40))
    add_accent_bar(slide, x, y, Inches(2.95), Inches(0.08), color)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.35), Inches(0.5), title, size=18, bold=True, color=color)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.9), Inches(2.35), Inches(1.5), desc, size=13, color=C_GRAY)

add_rounded_rect(slide, Inches(0.6), Inches(4.1), Inches(12.1), Inches(2.5), RGBColor(0x24, 0x2A, 0x40))
add_textbox(slide, Inches(1.0), Inches(4.3), Inches(11.3), Inches(0.5),
    "\U0001f6e1 資安與合規", size=18, bold=True, color=C_ACCENT)
security = [
    "數據加密：傳輸 TLS 1.3 + 靜態 AES-256 加密",
    "合規支援：符合 GDPR / CCPA / 台灣個資法規範",
    "部署彈性：可選公有雲（AWS/GCP/Azure）或私有雲 / 地端部署",
    "權限控管：RBAC 角色權限 + 操作審計日誌 + 資料脫敏處理",
]
for i, txt in enumerate(security):
    add_textbox(slide, Inches(1.2), Inches(4.9 + i * 0.4), Inches(11), Inches(0.4), f"\U000025b6 {txt}", size=13, color=C_GRAY)

# ════════════════════════════════════════════
# SLIDE 9: 導入效益
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "導入效益分析", "量化指標與預期成果")

metrics = [
    ("65%", "客服自動化比例", "一線問題由 AI Agent 自動解決", C_PRIMARY),
    ("3.2x", "行銷互動提升", "個人化推薦點擊率倍數成長", C_PURPLE),
    ("40%", "營運效率提升", "庫存管理與跟單自動化節省人力", C_GREEN),
    ("50%", "顧客滿意度提升", "即時回應 + 個人化體驗", C_ORANGE),
]

for i, (num, title, desc, color) in enumerate(metrics):
    x = Inches(0.6 + i * 3.15)
    y = Inches(1.3)
    card = add_rounded_rect(slide, x, y, Inches(2.95), Inches(2.8), RGBColor(0x24, 0x2A, 0x40))
    add_accent_bar(slide, x, y, Inches(2.95), Inches(0.08), color)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.55), Inches(1.0), num, size=48, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.4), Inches(2.55), Inches(0.5), title, size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.9), Inches(2.55), Inches(0.6), desc, size=12, color=C_GRAY, align=PP_ALIGN.CENTER)

# ROI estimate
add_rounded_rect(slide, Inches(0.6), Inches(4.4), Inches(12.1), Inches(2.2), RGBColor(0x24, 0x2A, 0x40))
add_textbox(slide, Inches(1.0), Inches(4.6), Inches(11.3), Inches(0.5),
    "\U0001f4b0 投資回報預估（以年營收 NT$5,000 萬之電商為例）", size=18, bold=True, color=C_ACCENT)
roi = [
    "人力節省：減少 4-6 名客服人員 → 年省約 NT$200-300 萬",
    "營收增長：購物車挽回 + 個人化推薦 → 年增營收 NT$500-800 萬",
    "IT 成本：方案費用約 NT$60-150 萬/年 → ROI 首年即可達 3-5 倍",
]
for i, txt in enumerate(roi):
    add_textbox(slide, Inches(1.2), Inches(5.2 + i * 0.4), Inches(11), Inches(0.4), f"\U000025b6 {txt}", size=14, color=C_GRAY)

# ════════════════════════════════════════════
# SLIDE 10: 合作方案
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "合作方案", "三種合作模式，彈性導入")

plans = [
    ("\U0001f331 輕量方案", "NT$ 60 萬 / 年", "• 單渠道 AI 客服機器人\n• 基本 CRM 整合\n• 5 個自動化流程\n• 每月成效報告\n• Email + LINE 支援", C_GREEN, "適合：小型電商\n快速導入"),
    ("\U0001f680 標準方案", "NT$ 120 萬 / 年", "• 多渠道 AI 客服 + 推薦\n• 完整 CRM 串接\n• 15 個自動化流程\n• A/B 測試 + 優化建議\n• 優先技術支援", C_PRIMARY, "適合：中型電商/零售\n全面導入"),
    ("\U0001f451 旗艦方案", "NT$ 240 萬 / 年", "• 全場景 AI Agent 部署\n• 私有化模型 + 地端/混合雲\n• 無限流程 + 客製開發\n• 專屬 AI 訓練與調校\n• 24/7 專屬客戶成功經理", C_PURPLE, "適合：大型零售集團\n深度定制"),
]

for i, (title, price, features, color, note) in enumerate(plans):
    x = Inches(0.5 + i * 4.2)
    card = add_rounded_rect(slide, x, Inches(1.3), Inches(3.9), Inches(4.8), RGBColor(0x24, 0x2A, 0x40))
    add_accent_bar(slide, x, Inches(1.3), Inches(3.9), Inches(0.08), color)
    add_textbox(slide, x + Inches(0.3), Inches(1.5), Inches(3.3), Inches(0.5), title, size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.3), Inches(2.0), Inches(3.3), Inches(0.5), price, size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.3), Inches(2.5), Inches(3.3), Inches(2.2), features, size=13, color=C_GRAY)
    add_textbox(slide, x + Inches(0.3), Inches(4.8), Inches(3.3), Inches(0.8), note, size=12, color=C_ACCENT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 11: 成功案例
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "成功案例參考", "Reference Cases")

cases = [
    ("某知名服飾電商", "導入 AI Agent 客服\n支援 LINE / FB / Web 三渠道\n客服量單月 12,000+ 則",
     "問題解決率 83%\n人力節省 58%\nNPS 由 32 提升至 68", C_PRIMARY),
    ("連鎖超市 (200+ 門市)", "店員 AI 助手\n智慧庫存管理\n會員辨識與推薦",
     "庫存成本降 22%\nVIP 客單提升 41%\n新訓時間縮短 60%", C_GREEN),
    ("跨境電商平台", "多語言 AI Agent\n24/7 自動客服\n自動跟單與物流查詢",
     "跨語言客服省 45% 成本\n購物車挽回率升至 31%\n年營收增 NT$1,200 萬", C_PURPLE),
]

for i, (title, desc, result, color) in enumerate(cases):
    x = Inches(0.5 + i * 4.2)
    card = add_rounded_rect(slide, x, Inches(1.3), Inches(3.9), Inches(3.5), RGBColor(0x24, 0x2A, 0x40))
    add_accent_bar(slide, x, Inches(1.3), Inches(3.9), Inches(0.08), color)
    add_textbox(slide, x + Inches(0.3), Inches(1.5), Inches(3.3), Inches(0.5), title, size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.3), Inches(2.1), Inches(3.3), Inches(1.0), desc, size=13, color=C_GRAY)
    add_textbox(slide, x + Inches(0.3), Inches(3.2), Inches(3.3), Inches(1.3), result, size=13, bold=True, color=C_ACCENT)

add_rounded_rect(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.5), RGBColor(0x24, 0x2A, 0x40))
add_textbox(slide, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.0),
    "\U0001f4ac 「導入 AI Agent 後，我們的客服團隊從被動接電話轉為主動經營會員關係，這是過去五年最有價值的技術投資。」\n— 某大型零售集團 IT 副總",
    size=14, color=C_GRAY)

# ════════════════════════════════════════════
# SLIDE 12: 結語
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), C_PRIMARY)
add_rect(slide, Inches(0), Inches(7.38), Inches(13.333), Inches(0.12), C_PRIMARY)

add_textbox(slide, Inches(1.2), Inches(1.5), Inches(11), Inches(0.6), "結語與下一步", size=18, color=C_ACCENT)
add_textbox(slide, Inches(1.2), Inches(2.2), Inches(11), Inches(1.2),
    "讓 AI Agent 成為您\nCRM 的智慧大腦", size=44, bold=True, color=C_WHITE)
add_accent_bar(slide, Inches(1.2), Inches(3.6), Inches(3), Inches(0.06), C_PRIMARY)

add_rounded_rect(slide, Inches(1.2), Inches(4.0), Inches(10.9), Inches(2.0), RGBColor(0x24, 0x2A, 0x40))
next_steps = [
    "\U0001f4ac Step 1：免費諮詢 — 了解貴公司現況與需求（約 1 小時）",
    "\U0001f4cb Step 2：POC 證明 — 選定一個場景進行 2 週概念驗證",
    "\U0001f680 Step 3：導入規劃 — 量身制定導入藍圖與時程",
    "\U0001f4c8 Step 4：上線優化 — 逐步上線 + 持續調校 + 成效追蹤",
]
for i, txt in enumerate(next_steps):
    add_textbox(slide, Inches(1.5), Inches(4.2 + i * 0.4), Inches(10.3), Inches(0.4), txt, size=16, color=C_GRAY)

add_textbox(slide, Inches(1.2), Inches(6.3), Inches(5), Inches(0.4),
    "聯絡我們：contact@aiagent-crm.com", size=14, color=C_GRAY)
add_textbox(slide, Inches(1.2), Inches(6.7), Inches(5), Inches(0.4),
    "官網：www.aiagent-crm.com", size=14, color=C_GRAY)

# decoration
circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(1.5), Inches(2.0), Inches(2.0))
circle2.fill.solid()
circle2.fill.fore_color.rgb = C_PRIMARY
circle2.fill.fore_color.brightness = 0.75
circle2.line.fill.background()

circle3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.5), Inches(0.8), Inches(1.2), Inches(1.2))
circle3.fill.solid()
circle3.fill.fore_color.rgb = C_ACCENT
circle3.fill.fore_color.brightness = 0.75
circle3.line.fill.background()

# ── Save ──
output_path = "AI_Agent_CRM_數位轉型方案.pptx"
prs.save(output_path)
print(f"✅ 簡報已生成: {output_path}")
