from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── light color palette ──
BG_PAGE = RGBColor(0xF1, 0xF5, 0xF9)
BG_CARD = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x0F, 0x17, 0x2A)
TEXT_BODY = RGBColor(0x33, 0x41, 0x55)
TEXT_SECONDARY = RGBColor(0x64, 0x74, 0x8B)
PRIMARY = RGBColor(0x25, 0x63, 0xEB)
ACCENT = RGBColor(0x0D, 0x94, 0x88)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
CODE_BG = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
CARD_SHADOW = RGBColor(0xF8, 0xFA, 0xFC)

FONT = "Microsoft JhengHei"
CODE_FONT = "Consolas"

def set_shape_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    set_shape_fill(shape, color)
    shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    set_shape_fill(shape, color)
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def add_accent_bar(slide, left, top, width, height, color):
    return add_rect(slide, left, top, width, height, color)

def section_header(slide, title, subtitle=""):
    add_bg(slide, BG_PAGE)
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(4.5), PRIMARY)
    add_textbox(slide, Inches(1), Inches(0.8), Inches(11), Inches(1.2), title, size=40, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(1), Inches(1.9), Inches(11), Inches(0.8), subtitle, size=22, color=RGBColor(0xBF, 0xDB, 0xFE))

# ════════════════════════════════════════════
# SLIDE 1: 封面
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(7.0), Inches(7.5), PRIMARY)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), PRIMARY)
add_rect(slide, Inches(0), Inches(7.4), Inches(13.333), Inches(0.1), PRIMARY)

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.6), "AI AGENT 快速開發方案", size=16, color=RGBColor(0xBF, 0xDB, 0xFE))

add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(2.0),
    "AI Agent\n快速打造小型\nWeb 會計系統", size=48, bold=True, color=WHITE)
add_accent_bar(slide, Inches(0.8), Inches(4.2), Inches(2), Inches(0.06), RGBColor(0xFD, 0xBA, 0x74))
add_textbox(slide, Inches(0.8), Inches(4.6), Inches(5.5), Inches(0.5),
    "從需求到部署 · 一週搞定一套會計系統", size=20, color=RGBColor(0xBF, 0xDB, 0xFE))
add_textbox(slide, Inches(0.8), Inches(5.3), Inches(5.5), Inches(0.5),
    "2026 年度技術方案提案", size=15, color=TEXT_SECONDARY)

add_textbox(slide, Inches(0.8), Inches(6.5), Inches(5.5), Inches(0.5),
    "Powered by AI Agent · Modern Web Stack", size=14, color=TEXT_SECONDARY)

# right side decoration
circle_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(1.5), Inches(4.0), Inches(4.0))
circle_bg.fill.solid()
circle_bg.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xFF)
circle_bg.line.fill.background()

add_textbox(slide, Inches(9.2), Inches(2.8), Inches(3.5), Inches(1.5),
    "⚡\n一週\n交付", size=36, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

deco_items = ["React + Tailwind", "FastAPI / Next.js", "PostgreSQL", "AI Agent 協作"]
for i, txt in enumerate(deco_items):
    y = Inches(5.0 + i * 0.45)
    add_rounded_rect(slide, Inches(8.5), y, Inches(4.0), Inches(0.35), WHITE, BORDER)
    add_textbox(slide, Inches(8.7), y + Inches(0.02), Inches(3.6), Inches(0.3), f"  \U000025b6 {txt}", size=13, color=TEXT_BODY)

# ════════════════════════════════════════════
# SLIDE 2: 議程
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "議 程 內 容", "AGENDA")

agenda_items = [
    ("01", "傳統會計開發痛點", "為什麼需要 AI Agent"),
    ("02", "AI Agent 的開發角色", "貫穿開發全流程的 AI 協作"),
    ("03", "系統功能規劃", "小型會計系統必備模組"),
    ("04", "系統架構設計", "Frontend + Backend + DB + AI"),
    ("05", "AI Agent 開發流程", "Prompt → 生成 → Review → 迭代"),
    ("06", "核心程式碼展示", "關鍵實作範例與 Prompt"),
    ("07", "部署方案", "Docker / Vercel + Supabase"),
    ("08", "導入效益與結語", "量化效益與下一步行動"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(2.6 + row * 1.2)

    card = add_rounded_rect(slide, x, y, Inches(5.8), Inches(1.0), BG_CARD, BORDER)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.05), Inches(0.7), Inches(0.9), num, size=24, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(1.0), y + Inches(0.05), Inches(4.5), Inches(0.45), title, size=18, bold=True, color=TEXT_DARK)
    add_textbox(slide, x + Inches(1.0), y + Inches(0.5), Inches(4.5), Inches(0.4), desc, size=15, color=TEXT_SECONDARY)

# ════════════════════════════════════════════
# SLIDE 3: 傳統會計開發痛點
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "傳統會計開發痛點", "小型會計系統為何需要 AI Agent？")

pain_points = [
    ("開發週期長", "從需求釐清到上線\n平均需要 3~6 個月\n規格確認反覆耗時", PRIMARY),
    ("規格變動頻", "會計準則與流程經常調整\n需求變更導致大量重工\n模組耦合難以快速修改", ORANGE),
    ("測試耗時", "手動測試佔開發時間 40%\n邊界案例難以窮舉\n財務數值錯誤風險高", PURPLE),
    ("維護成本高", "Bug fix + 功能迭代不間斷\n缺乏文件導致交接困難\n技術債隨時間累積", GREEN),
]

for i, (title, desc, color) in enumerate(pain_points):
    x = Inches(0.6 + i * 3.15)
    y = Inches(2.6)
    card = add_rounded_rect(slide, x, y, Inches(2.95), Inches(2.8), BG_CARD, BORDER)
    add_accent_bar(slide, x, y, Inches(2.95), Inches(0.06), color)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.35), Inches(0.5), title, size=22, bold=True, color=color)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.9), Inches(2.35), Inches(1.8), desc, size=16, color=TEXT_BODY)

add_rounded_rect(slide, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.2), BG_CARD, BORDER)
add_textbox(slide, Inches(1.0), Inches(5.9), Inches(11.3), Inches(0.7),
    "\U0001f4a1 核心問題：傳統開發模式無法快速因應小型企業的預算限制與迭代需求，AI Agent 能將開發週期從「月」縮短到「天」",
    size=17, bold=False, color=TEXT_BODY)

# ════════════════════════════════════════════
# SLIDE 4: AI Agent 的開發角色
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "AI Agent 的開發角色", "貫穿開發全流程的 AI 協作")

roles = [
    ("\U0001f4dd 程式碼生成", "根據需求描述直接生成 CRUD\n從 UI Component 到 API Route\n資料庫 Schema 自動設計", PRIMARY),
    ("\U0001f50d 除錯與優化", "錯誤訊息自動分析與修復\n效能瓶頸檢測與建議\n重構與最佳實踐導入", PURPLE),
    ("\U0001f4ca 測試自動化", "單元測試、整合測試自動生成\n邊界案例自動補全\n測試覆蓋率可達 90%+", GREEN),
    ("\U0001f4c4 文件生成", "API 文件自動產生\n資料庫 ER 圖自動繪製\n部署與使用說明自動撰寫", ORANGE),
]

for i, (title, desc, color) in enumerate(roles):
    x = Inches(0.6 + i * 3.15)
    y = Inches(2.6)
    card = add_rounded_rect(slide, x, y, Inches(2.95), Inches(2.8), BG_CARD, BORDER)
    add_accent_bar(slide, x, y, Inches(2.95), Inches(0.06), color)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.35), Inches(0.5), title, size=20, bold=True, color=color)
    add_textbox(slide, x + Inches(0.3), y + Inches(1.0), Inches(2.35), Inches(1.6), desc, size=16, color=TEXT_BODY)

add_rounded_rect(slide, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.2), BG_CARD, BORDER)
add_textbox(slide, Inches(1.0), Inches(5.9), Inches(11.3), Inches(0.7),
    "\U0001f916 AI Agent 不是取代開發者，而是讓開發者專注在高階設計與業務邏輯，重複性勞動交給 AI",
    size=17, bold=False, color=TEXT_BODY)

# ════════════════════════════════════════════
# SLIDE 5: 系統功能規劃
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "系統功能規劃", "小型會計系統必備五大模組")

modules = [
    ("\U0001f4cb 憑證管理", "• 傳票新增 / 編輯 / 查詢\n• 複式記帳自動分錄\n• 支援多幣別\n• 附件上傳與管理", PRIMARY),
    ("\U0001f4ca 總帳管理", "• 會計科目維護\n• 科目餘額表\n• 明細分類帳\n• 試算表自動產生", PURPLE),
    ("\U0001f4b3 應收/應付", "• 客戶/供應商管理\n• 帳齡分析\n• 收款/付款預警\n• 沖帳作業", GREEN),
    ("\U0001f4c8 財務報表", "• 損益表\n• 資產負債表\n• 現金流量表\n• 匯出 Excel / PDF", ORANGE),
    ("\U0001f464 系統管理", "• 使用者與權限管理\n• 操作審計日誌\n• 系統參數設定\n• 資料備份與還原", ACCENT),
]

for i, (title, desc, color) in enumerate(modules):
    x = Inches(0.4 + i * 2.55)
    y = Inches(2.6)
    card = add_rounded_rect(slide, x, y, Inches(2.35), Inches(3.5), BG_CARD, BORDER)
    add_accent_bar(slide, x, y, Inches(2.35), Inches(0.06), color)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.3), Inches(1.95), Inches(0.5), title, size=18, bold=True, color=color)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.9), Inches(1.95), Inches(2.4), desc, size=14, color=TEXT_BODY)

add_rounded_rect(slide, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.8), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.5),
    "以上功能可用 AI Agent 在 3~5 天內完成 CRUD 開發，透過迭代逐步完善", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 6: 系統架構
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "系統架構設計", "Frontend + Backend + DB + AI Agent")

layers = [
    ("\U0001f310 前端層 Frontend", "React + TypeScript + Tailwind CSS\nReact Router 路由管理 · Shadcn/ui 元件庫\nVite 建置工具 · Axios API 串接", PRIMARY),
    ("\U0001f4e1 API 層 Backend", "Python FastAPI / Node.js Next.js\nRESTful API 設計 · JWT 認證\nSQLAlchemy ORM · Pydantic 驗證", PURPLE),
    ("\U0001f4be 資料層 Database", "PostgreSQL / SQLite（開發用）\nAlembic 資料庫遷移\n索引優化 · 交易事務管理", GREEN),
    ("\U0001f916 AI Agent 協作", "Cursor IDE · Claude / GPT-4\nPrompt Engineering 驅動開發\n版本控制 + AI Review", ORANGE),
]

for i, (title, desc, color) in enumerate(layers):
    y = Inches(2.5 + i * 1.2)
    card = add_rounded_rect(slide, Inches(0.8), y, Inches(11.7), Inches(1.05), BG_CARD, BORDER)
    add_accent_bar(slide, Inches(0.8), y, Inches(0.08), Inches(1.05), color)
    add_textbox(slide, Inches(1.2), y + Inches(0.05), Inches(3), Inches(0.45), title, size=20, bold=True, color=color)
    add_textbox(slide, Inches(1.2), y + Inches(0.5), Inches(10.5), Inches(0.5), desc, size=15, color=TEXT_BODY)

add_rounded_rect(slide, Inches(0.8), Inches(7.0 - 0.4), Inches(11.7), Inches(0.35), RGBColor(0xEE, 0xF2, 0xFF), BORDER)
add_textbox(slide, Inches(1.2), Inches(7.0 - 0.37), Inches(11.3), Inches(0.3),
    "\U0001f517 資料流：使用者操作 → React UI → API 請求 → FastAPI → ORM → PostgreSQL → 回應渲染", size=14, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 7: AI Agent 開發流程
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "AI Agent 開發流程", "Prompt → 生成 → Review → 迭代")

steps = [
    ("Step 1", "需求描述", "以自然語言描述功能需求\n明確輸入/輸出/商業邏輯\n附上參考 UI 或資料範例", PRIMARY),
    ("Step 2", "AI 生成", "Agent 產出完整 CRUD 程式碼\n包含 Model + API + UI\n資料庫 Migration 腳本", PURPLE),
    ("Step 3", "人工 Review", "檢查業務邏輯正確性\n調整 UI/UX 細節\n安全性與權限驗證", ORANGE),
    ("Step 4", "迭代優化", "提出修改需求給 Agent\n自動更新相關聯程式碼\n回歸測試確認無影響", GREEN),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.6 + i * 3.15)
    y = Inches(2.6)
    card = add_rounded_rect(slide, x, y, Inches(2.95), Inches(2.8), BG_CARD, BORDER)
    add_accent_bar(slide, x, y, Inches(2.95), Inches(0.06), color)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.1), Inches(2.55), Inches(0.4), num, size=14, bold=True, color=color)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.4), Inches(2.55), Inches(0.5), title, size=22, bold=True, color=TEXT_DARK)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.0), Inches(2.55), Inches(1.6), desc, size=16, color=TEXT_BODY)

# Prompt example
add_rounded_rect(slide, Inches(0.6), Inches(5.7), Inches(5.5), Inches(1.4), CODE_BG)
add_textbox(slide, Inches(0.9), Inches(5.8), Inches(5.0), Inches(0.3), "  範例 Prompt", size=14, bold=True, color=RGBColor(0x82, 0xCA, 0x9D), font_name=CODE_FONT)
prompt_text = (
    '  "請幫我建立一個傳票輸入頁面，包含：\n'
    '  日期選擇、科目下拉選單、借方/貸方金額、\n'
    '  摘要文字欄位。使用 React + TypeScript，\n'
    '  表單驗證使用 react-hook-form。"'
)
add_textbox(slide, Inches(0.9), Inches(6.15), Inches(5.0), Inches(0.9), prompt_text, size=13, color=RGBColor(0xE2, 0xE8, 0xF0), font_name=CODE_FONT)

add_rounded_rect(slide, Inches(6.4), Inches(5.7), Inches(6.3), Inches(1.4), BG_CARD, BORDER)
add_textbox(slide, Inches(6.7), Inches(5.85), Inches(5.7), Inches(0.3), "  \U0001f916 AI Agent 回應", size=14, bold=True, color=TEXT_DARK)
add_textbox(slide, Inches(6.7), Inches(6.2), Inches(5.7), Inches(0.7),
    "  自動產出完整的 Component 程式碼\n  含表單驗證、API 串接、TypeScript 型別\n  以及對應的 Backend API Route", size=14, color=TEXT_BODY)

# ════════════════════════════════════════════
# SLIDE 8: 核心程式碼展示
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "核心程式碼展示", "關鍵實作範例與 AI 生成對比")

# Code example 1
add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.8), Inches(0.4),
    "\U0001f4be 會計分錄 Entity 定義（Python / SQLAlchemy）", size=18, bold=True, color=TEXT_DARK)

code1 = (
    '  class JournalEntry(Base):\n'
    '      __tablename__ = "journal_entries"\n'
    '      id = Column(UUID, primary_key=True)\n'
    '      date = Column(Date, nullable=False)\n'
    '      description = Column(String(500))\n'
    '      status = Column(Enum(EntryStatus))\n'
    '  \n'
    '  class JournalLine(Base):\n'
    '      __tablename__ = "journal_lines"\n'
    '      id = Column(UUID, primary_key=True)\n'
    '      entry_id = Column(ForeignKey("journal_entries.id"))\n'
    '      account_id = Column(ForeignKey("accounts.id"))\n'
    '      debit = Column(Numeric(14,2), default=0)\n'
    '      credit = Column(Numeric(14,2), default=0)'
)
add_rounded_rect(slide, Inches(0.6), Inches(3.0), Inches(5.8), Inches(2.5), CODE_BG)
add_textbox(slide, Inches(0.8), Inches(3.1), Inches(5.4), Inches(2.3), code1, size=13, color=RGBColor(0xE2, 0xE8, 0xF0), font_name=CODE_FONT)

# Code example 2
add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.8), Inches(0.4),
    "\U0001f4b0 複式記帳驗證邏輯", size=18, bold=True, color=TEXT_DARK)

code2 = (
    '  def validate_double_entry(lines: list[Line]) -> bool:\n'
    '      total_debit = sum(l.debit for l in lines)\n'
    '      total_credit = sum(l.credit for l in lines)\n'
    '      if total_debit != total_credit:\n'
    '          raise ValueError(\n'
    '              f"借貸不平衡: 借方{total_debit} "\n'
    '              f"≠ 貸方{total_credit}"\n'
    '          )\n'
    '      return True\n'
    '  \n'
    '  # AI Agent 自動生成的測試案例\n'
    '  def test_validate_double_entry():\n'
    '      lines = [\n'
    '          Line(debit=1000, credit=0),\n'
    '          Line(debit=0, credit=1000),\n'
    '      ]\n'
    '      assert validate_double_entry(lines)'
)
add_rounded_rect(slide, Inches(6.8), Inches(3.0), Inches(5.8), Inches(2.5), CODE_BG)
add_textbox(slide, Inches(7.0), Inches(3.1), Inches(5.4), Inches(2.3), code2, size=13, color=RGBColor(0xE2, 0xE8, 0xF0), font_name=CODE_FONT)

add_rounded_rect(slide, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.2), BG_CARD, BORDER)
add_textbox(slide, Inches(1.0), Inches(5.9), Inches(11.3), Inches(0.7),
    "\U0001f916 以上程式碼 80% 由 AI Agent 生成，開發者只需描述業務邏輯與資料結構，Agent 即可產出完整 Model、API、UI 及測試",
    size=16, color=TEXT_BODY)

# ════════════════════════════════════════════
# SLIDE 9: 資料安全與合規
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "資料安全與合規", "會計系統的資安考量")

sec_items = [
    ("\U0001f512 資料加密", "傳輸加密 TLS 1.3\n靜態資料 AES-256 加密\n敏感欄位脫敏處理", PRIMARY),
    ("\U0001f6e1 存取控制", "RBAC 角色權限管理\nJWT Token 認證\n操作審計日誌完整記錄", PURPLE),
    ("\U0001f4cb 資料備份", "自動每日備份\n增量備份機制\n異地備援選項", GREEN),
    ("\U00002699 法規合規", "符合一般會計原則\n審計軌跡保留\n支援匯出供會計師查核", ORANGE),
]

for i, (title, desc, color) in enumerate(sec_items):
    x = Inches(0.6 + i * 3.15)
    y = Inches(2.6)
    card = add_rounded_rect(slide, x, y, Inches(2.95), Inches(2.6), BG_CARD, BORDER)
    add_accent_bar(slide, x, y, Inches(2.95), Inches(0.06), color)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.35), Inches(0.5), title, size=20, bold=True, color=color)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.9), Inches(2.35), Inches(1.5), desc, size=15, color=TEXT_BODY)

add_rounded_rect(slide, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.3), BG_CARD, BORDER)
add_textbox(slide, Inches(1.0), Inches(5.7), Inches(11.3), Inches(0.7),
    "\U0001f4a1 小型會計系統採用雲端部署，安全等級等同一般 SaaS 產品。若需地端部署，可將整個系統打包成 Docker Image 佈署於客戶機房",
    size=16, color=TEXT_BODY)

# ════════════════════════════════════════════
# SLIDE 10: 部署方案
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "部署方案", "彈性部署選擇，滿足不同需求")

deployments = [
    ("\U0001f680 快速啟動", "Vercel + Supabase\n• 前端部署 Vercel (Free Tier)\n• 資料庫用 Supabase (Free Tier)\n• 後端可選 Vercel Serverless\n• 適合 POC 與小型團隊\n• 零成本啟動", PRIMARY),
    ("\U0001f4e6 Docker 部署", "Docker Compose 一鍵啟動\n• Frontend + Backend + DB\n• Dockerfile + docker-compose.yml\n• 支援任何雲端主機\n• 適合正式環境\n• 可攜帶至地端", PURPLE),
    ("\U0001f3e2 混合方案", "Vercel + 後端主機\n• 前端 Vercel 全球 CDN\n• 後端佈署於雲端 VM\n• 資料庫可選 RDS / Supabase\n• 兼顧效能與成本\n• 適合中型團隊", GREEN),
]

for i, (title, desc, color) in enumerate(deployments):
    x = Inches(0.5 + i * 4.2)
    card = add_rounded_rect(slide, x, Inches(2.6), Inches(3.9), Inches(3.5), BG_CARD, BORDER)
    add_accent_bar(slide, x, Inches(2.6), Inches(3.9), Inches(0.06), color)
    add_textbox(slide, x + Inches(0.3), Inches(2.8), Inches(3.3), Inches(0.5), title, size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.3), Inches(3.4), Inches(3.3), Inches(2.5), desc, size=15, color=TEXT_BODY)

add_rounded_rect(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8), PRIMARY)
add_textbox(slide, Inches(0.9), Inches(6.45), Inches(11.7), Inches(0.5),
    "\U0001f916 部署 Dockerfile 與 CI/CD Pipeline 也可由 AI Agent 自動生成，只需描述部署環境即可",
    size=16, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 11: 導入效益
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "導入效益", "AI Agent 開發 vs 傳統開發對比")

metrics = [
    ("開發時間", "傳統 3~6 個月\nAI Agent 1~2 週\n縮短 80~90%", PRIMARY),
    ("開發成本", "傳統 NT$80~150 萬\nAI Agent NT$15~30 萬\n節省 70~80%", PURPLE),
    ("迭代速度", "傳統 2~4 週/次\nAI Agent 2~3 天/次\n提升 5~10 倍", GREEN),
    ("測試覆蓋", "傳統 50~60%\nAI Agent 85~95%\n提升 40%", ORANGE),
]

for i, (title, trad_val, color) in enumerate(metrics):
    x = Inches(0.6 + i * 3.15)
    y = Inches(2.6)
    card = add_rounded_rect(slide, x, y, Inches(2.95), Inches(3.8), BG_CARD, BORDER)
    add_accent_bar(slide, x, y, Inches(2.95), Inches(0.06), color)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.55), Inches(0.5), title, size=22, bold=True, color=color, align=PP_ALIGN.CENTER)

    # Traditional
    lines = trad_val.split("\n")
    add_textbox(slide, x + Inches(0.2), y + Inches(1.0), Inches(2.55), Inches(1.0), lines[0], size=14, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)

    # Arrow
    add_textbox(slide, x + Inches(0.2), y + Inches(1.6), Inches(2.55), Inches(0.5), "  \U000025bc", size=16, color=color, align=PP_ALIGN.CENTER)

    # AI Agent value
    add_textbox(slide, x + Inches(2.55 - 1.0), y + Inches(2.0), Inches(2.0), Inches(0.5), lines[1], size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), y + Inches(2.5), Inches(2.55), Inches(0.4), lines[2], size=15, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # Highlight box
    highlight = add_rounded_rect(slide, x + Inches(0.4), y + Inches(3.0), Inches(2.15), Inches(0.6), color)
    highlight.fill.fore_color.brightness = 0.85
    if len(lines) >= 3:
        add_textbox(slide, x + Inches(0.4), y + Inches(3.1), Inches(2.15), Inches(0.4), lines[2], size=14, bold=True, color=color, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 12: 結語
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(7.0), Inches(7.5), PRIMARY)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), PRIMARY)
add_rect(slide, Inches(0), Inches(7.4), Inches(13.333), Inches(0.1), PRIMARY)

add_textbox(slide, Inches(0.8), Inches(1.0), Inches(5.5), Inches(0.6), "結語與下一步", size=16, color=RGBColor(0xBF, 0xDB, 0xFE))
add_textbox(slide, Inches(0.8), Inches(1.8), Inches(5.8), Inches(1.5),
    "AI Agent 讓會計系統開發\n不再是大型團隊的專利", size=40, bold=True, color=WHITE)
add_accent_bar(slide, Inches(0.8), Inches(3.5), Inches(2), Inches(0.06), RGBColor(0xFD, 0xBA, 0x74))
add_textbox(slide, Inches(0.8), Inches(3.9), Inches(5.5), Inches(0.5),
    "一個人 + AI Agent 一週即可交付", size=20, color=RGBColor(0xBF, 0xDB, 0xFE))

add_rounded_rect(slide, Inches(0.8), Inches(4.6), Inches(5.5), Inches(2.2), RGBColor(0x1E, 0x3A, 0x5F))
next_steps = [
    "\U0001f4ac Step 1：需求訪談 — 了解目前會計作業流程",
    "\U0001f4cb Step 2：2 週 POC — 建立核心憑證 + 報表模組",
    "\U0001f680 Step 3：系統交付 — 完整功能 + 教育訓練",
    "\U0001f4c8 Step 4：持續迭代 — 依使用反饋持續優化",
]
for i, txt in enumerate(next_steps):
    add_textbox(slide, Inches(1.1), Inches(4.8 + i * 0.45), Inches(5.0), Inches(0.4), txt, size=15, color=WHITE)

add_textbox(slide, Inches(0.8), Inches(6.9), Inches(5.5), Inches(0.4),
    "聯絡我們：dev@aiagent-accounting.com", size=13, color=TEXT_SECONDARY)

# right side
circle_r = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(1.5), Inches(3.5), Inches(3.5))
circle_r.fill.solid()
circle_r.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xFF)
circle_r.line.fill.background()

add_textbox(slide, Inches(8.8), Inches(2.5), Inches(3.0), Inches(1.5),
    "馬上開始\n\U0001f680", size=36, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

items_r = ["\U0001f4e9 填寫需求表單", "\U0001f4ac 預約線上 Demo", "\U0001f916 立即體驗 POC"]
for i, txt in enumerate(items_r):
    add_rounded_rect(slide, Inches(8.5), Inches(5.2 + i * 0.5), Inches(4.0), Inches(0.4), WHITE, BORDER)
    add_textbox(slide, Inches(8.7), Inches(5.25 + i * 0.5), Inches(3.6), Inches(0.35), txt, size=15, color=TEXT_BODY)

# ── Save ──
output_path = "AI_Agent_會計系統快速開發方案.pptx"
prs.save(output_path)
print(f"✅ 簡報已生成: {output_path}")
