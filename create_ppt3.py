from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

C_DARK = RGBColor(0x0F, 0x0B, 0x1A)
C_CARD = RGBColor(0x1A, 0x15, 0x2E)
C_PRIMARY = RGBColor(0x81, 0x63, 0xFC)
C_ACCENT = RGBColor(0x2D, 0xD4, 0xBF)
C_PINK = RGBColor(0xF4, 0x72, 0xB6)
C_ORANGE = RGBColor(0xFB, 0x92, 0x3C)
C_GREEN = RGBColor(0x34, 0xD3, 0x99)
C_YELLOW = RGBColor(0xFC, 0xD3, 0x4F)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
C_SECONDARY = RGBColor(0x88, 0x88, 0xAA)
C_SUBTLE = RGBColor(0x2A, 0x25, 0x45)
FONT = "Microsoft JhengHei"

def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    return s

def rrect(slide, l, t, w, h, c, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c
    if border: s.line.color.rgb = border; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, size=18, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT, font=FONT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = font; p.alignment = align
    return tf

def header(slide, title, sub=""):
    add_bg(slide, C_DARK)
    rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), C_PRIMARY)
    rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), C_PRIMARY)
    tb(slide, Inches(0.9), Inches(0.5), Inches(11), Inches(1.0), title, 38, True, C_WHITE)
    if sub:
        tb(slide, Inches(0.9), Inches(1.4), Inches(11), Inches(0.6), sub, 20, color=C_SECONDARY)

# ════════════════════════════════════════════
# SLIDE 1: 封面
# ════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, C_DARK)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.08), C_PRIMARY)
rect(sl, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), C_PRIMARY)

# gradient circles
for i, (x, r, c) in enumerate([
    (Inches(9.0), Inches(3.5), C_PRIMARY),
    (Inches(10.2), Inches(2.5), C_ACCENT),
    (Inches(7.8), Inches(2.0), C_PINK),
]):
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(0.8 + i*0.6), r, r)
    circ.fill.solid(); circ.fill.fore_color.rgb = c
    circ.fill.fore_color.brightness = 0.65; circ.line.fill.background()

tb(sl, Inches(1.0), Inches(0.8), Inches(11), Inches(0.5), "AI AGENT 技術深究系列", 16, color=C_SECONDARY)
tb(sl, Inches(1.0), Inches(1.6), Inches(7), Inches(1.5), "AI Agent 常用\nSkills 介紹", 50, True, C_WHITE)
rect(sl, Inches(1.0), Inches(3.3), Inches(2.5), Inches(0.06), C_PRIMARY)
tb(sl, Inches(1.0), Inches(3.7), Inches(7), Inches(0.5), "從工具調用到多模態理解，全面解析 AI Agent 的核心技能組", 18, color=C_SECONDARY)
tb(sl, Inches(1.0), Inches(4.5), Inches(5), Inches(0.4), "2026 技術專題", 14, color=C_SECONDARY)
tb(sl, Inches(1.0), Inches(5.0), Inches(5), Inches(0.4), "Powered by AI Agent Ecosystem", 13, color=C_SECONDARY)

# skill tags on right
tags = ["\U0001f50d Web Search", "\U0001f4bb Code Gen", "\U0001f4c8 Analysis",
        "\U0001f5bc Image", "\U0001f4e1 API", "\U0001f9e0 Planning",
        "\U0001f4be Memory", "\U0001f310 Multi-modal"]
for i, tag in enumerate(tags):
    row, col = i // 2, i % 2
    x, y = Inches(8.5 + col * 2.3), Inches(5.5 + row * 0.45)
    rrect(sl, x, y, Inches(2.1), Inches(0.38), C_SUBTLE, C_SUBTLE)
    tb(sl, x + Inches(0.15), y + Inches(0.02), Inches(1.8), Inches(0.34), tag, 12, color=C_GRAY)

# ════════════════════════════════════════════
# SLIDE 2: 議程
# ════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "議 程 內 容", "AGENDA")

items = [
    ("01", "什麼是 AI Agent Skills？", "技能系統的核心概念"),
    ("02", "資訊獲取類 Skills", "Web Search · 檔案讀取 · 網頁爬取"),
    ("03", "程式執行類 Skills", "Code Gen · Code Run · Debug"),
    ("04", "內容生成類 Skills", "文字 · 圖片 · 文件生成"),
    ("05", "資料處理類 Skills", "分析 · 視覺化 · 資料庫操作"),
    ("06", "工具串接類 Skills", "API · Function Calling · 通訊"),
    ("07", "進階 Skills", "規劃 · 記憶 · 多模態理解"),
    ("08", "實戰組合應用", "如何組合 Skills 解決真實問題"),
]

for i, (n, t, d) in enumerate(items):
    r, c = i // 2, i % 2
    x, y = Inches(0.8 + c * 6.2), Inches(2.4 + r * 1.25)
    rrect(sl, x, y, Inches(5.8), Inches(1.05), C_CARD, C_SUBTLE)
    tb(sl, x + Inches(0.25), y + Inches(0.05), Inches(0.7), Inches(0.95), n, 26, True, C_PRIMARY, PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.95), y + Inches(0.05), Inches(4.6), Inches(0.45), t, 20, True, C_WHITE)
    tb(sl, x + Inches(0.95), y + Inches(0.55), Inches(4.6), Inches(0.4), d, 15, color=C_SECONDARY)

# ════════════════════════════════════════════
# SLIDE 3: 什麼是 AI Agent Skills？
# ════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "什麼是 AI Agent Skills？", "技能系統的核心概念")

# definition box
rrect(sl, Inches(0.6), Inches(2.3), Inches(12.1), Inches(1.5), C_CARD, C_SUBTLE)
rect(sl, Inches(0.6), Inches(2.3), Inches(0.08), Inches(1.5), C_PRIMARY)
tb(sl, Inches(1.0), Inches(2.4), Inches(11.3), Inches(0.4), "\U0001f4a1 什麼是 Skill？", 22, True, C_PRIMARY)
tb(sl, Inches(1.0), Inches(2.9), Inches(11.3), Inches(0.7),
    "Skill 是 AI Agent 可以調用的具體能力單元，類似人類的「工具使用技能」。每個 Skill 封裝了特定的功能介面，\nAgent 根據任務需求動態選擇並組合多個 Skills 來完成複雜任務。", 16, color=C_GRAY)

# 3 key concepts
concepts = [
    ("\U0001f9f0 模組化", "每個 Skill 獨立封裝\n可插拔、可替換\n降低系統耦合度", C_PRIMARY),
    ("\U0001f517 可組合", "多個 Skills 可串聯使用\n形成複雜工作流\nAgent 自動編排順序", C_ACCENT),
    ("\U0001f4a1 可擴展", "可自訂開發新 Skill\n社群生態持續增長\n依場景動態載入", C_PINK),
]

for i, (t, d, c) in enumerate(concepts):
    x = Inches(0.6 + i * 4.2)
    rrect(sl, x, Inches(4.1), Inches(3.9), Inches(2.5), C_CARD, C_SUBTLE)
    rect(sl, x, Inches(4.1), Inches(3.9), Inches(0.06), c)
    tb(sl, x + Inches(0.3), Inches(4.3), Inches(3.3), Inches(0.5), t, 22, True, c)
    tb(sl, x + Inches(0.3), Inches(5.0), Inches(3.3), Inches(1.4), d, 16, color=C_GRAY)

# ════════════════════════════════════════════
# SLIDE 4: 核心技能分類總覽
# ════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "核心技能分類總覽", "五大分類 · 常見 Skills 一覽")

categories = [
    ("\U0001f50d 資訊獲取", "Web Search\nFile Reader\nWeb Scraper\nRSS Reader", C_PRIMARY),
    ("\U0001f4bb 程式執行", "Code Generator\nCode Interpreter\nDebugger\nShell Exec", C_ACCENT),
    ("\U0001f3a8 內容生成", "Text Generator\nImage Generator\nDocument Maker\nSlide Creator", C_PINK),
    ("\U0001f4ca 資料處理", "Data Analyzer\nChart Generator\nDB Query\nExcel Processor", C_YELLOW),
    ("\U0001f4e1 工具串接", "API Caller\nFunction Calling\nEmail Sender\nWebhook", C_GREEN),
]

for i, (cat, skills, color) in enumerate(categories):
    x = Inches(0.4 + i * 2.55)
    rrect(sl, x, Inches(2.4), Inches(2.35), Inches(4.2), C_CARD, C_SUBTLE)
    rect(sl, x, Inches(2.4), Inches(2.35), Inches(0.06), color)
    tb(sl, x + Inches(0.2), Inches(2.6), Inches(1.95), Inches(0.5), cat, 20, True, color, PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.2), Inches(3.2), Inches(1.95), Inches(3.0), skills, 15, color=C_GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 5: 資訊獲取類 Skills
# ════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "資訊獲取類 Skills", "Web Search · 檔案讀取 · 網頁爬取")

skills_detail = [
    ("\U0001f50d Web Search", "• 即時搜尋網際網路資訊\n• 支援 Google / Bing / DuckDuckGo\n• 回結構化搜尋結果\n• 可指定搜尋數量與來源", C_PRIMARY),
    ("\U0001f4c4 File Reader", "• 讀取本機與遠端檔案\n• 支援 PDF / Word / Excel / TXT\n• 自動解析文件結構\n• 搭配 OCR 辨識掃描檔", C_ACCENT),
    ("\U0001f578 Web Scraper", "• 提取網頁結構化資料\n• 支援 CSS Selector / XPath\n• JavaScript 渲染頁面\n• 遵守 robots.txt 規範", C_PINK),
]

for i, (t, d, c) in enumerate(skills_detail):
    x = Inches(0.5 + i * 4.2)
    rrect(sl, x, Inches(2.4), Inches(3.9), Inches(2.5), C_CARD, C_SUBTLE)
    rect(sl, x, Inches(2.4), Inches(3.9), Inches(0.06), c)
    tb(sl, x + Inches(0.3), Inches(2.6), Inches(3.3), Inches(0.5), t, 22, True, c)
    tb(sl, x + Inches(0.3), Inches(3.2), Inches(3.3), Inches(1.5), d, 16, color=C_GRAY)

# use case
rrect(sl, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.5), C_CARD, C_SUBTLE)
rect(sl, Inches(0.5), Inches(5.2), Inches(0.08), Inches(1.5), C_PRIMARY)
tb(sl, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.4), "\U0001f4a1 應用場景", 20, True, C_PRIMARY)
tb(sl, Inches(0.9), Inches(5.8), Inches(11.5), Inches(0.7),
    "市場競品分析：Web Search 蒐集競品資訊 → Web Scraper 提取結構化數據 → File Reader 讀取既有報告 → Agent 綜合產出分析簡報",
    16, color=C_GRAY)

# ════════════════════════════════════════════
# SLIDE 6: 程式執行類 Skills
# ════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "程式執行類 Skills", "Code Gen · Code Run · Debug")

skills_detail2 = [
    ("\U0001f4bb Code Generator", "• 根據需求描述生成程式碼\n• 支援多語言（Python/JS/TS/Go）\n• 含型別定義與錯誤處理\n• 自動產出單元測試", C_PRIMARY),
    ("\U0001f680 Code Interpreter", "• 隔離沙箱執行 Python 程式碼\n• 即時執行並回傳結果\n• 支援套件安裝\n• 資料分析與圖表繪製", C_ACCENT),
    ("\U0001f41b Debugger", "• 自動分析錯誤堆疊\n• 定位 Bug 位置與原因\n• 提供修復建議\n• 支援多語言除錯", C_PINK),
]

for i, (t, d, c) in enumerate(skills_detail2):
    x = Inches(0.5 + i * 4.2)
    rrect(sl, x, Inches(2.4), Inches(3.9), Inches(2.5), C_CARD, C_SUBTLE)
    rect(sl, x, Inches(2.4), Inches(3.9), Inches(0.06), c)
    tb(sl, x + Inches(0.3), Inches(2.6), Inches(3.3), Inches(0.5), t, 22, True, c)
    tb(sl, x + Inches(0.3), Inches(3.2), Inches(3.3), Inches(1.5), d, 16, color=C_GRAY)

rrect(sl, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.5), C_CARD, C_SUBTLE)
rect(sl, Inches(0.5), Inches(5.2), Inches(0.08), Inches(1.5), C_ACCENT)
tb(sl, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.4), "\U0001f4a1 應用場景", 20, True, C_ACCENT)
tb(sl, Inches(0.9), Inches(5.8), Inches(11.5), Inches(0.7),
    "數據分析自動化：Code Generator 產出分析腳本 → Code Interpreter 執行並產出圖表 → Debugger 自動修正錯誤 → 輸出最終報告",
    16, color=C_GRAY)

# ════════════════════════════════════════════
# SLIDE 7: 內容生成類 Skills
# ════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "內容生成類 Skills", "文字 · 圖片 · 文件生成")

skills_detail3 = [
    ("\U0001f4dd Text Generator", "• 文章/報告/郵件自動撰寫\n• 多語言翻譯與在地化\n• 摘要生成與改寫\n• SEO 優化內容產出", C_PRIMARY),
    ("\U0001f5bc Image Generator", "• 根據描述生成圖片\n• DALL·E / Stable Diffusion\n• 圖片編輯與風格轉換\n• 商業素材快速產出", C_ACCENT),
    ("\U0001f4c4 Document Maker", "• 自動生成 PDF / PPTX / DOCX\n• 報表與簡報自動排版\n• 資料驅動的模板填充\n• 批量文件生成", C_PINK),
]

for i, (t, d, c) in enumerate(skills_detail3):
    x = Inches(0.5 + i * 4.2)
    rrect(sl, x, Inches(2.4), Inches(3.9), Inches(2.5), C_CARD, C_SUBTLE)
    rect(sl, x, Inches(2.4), Inches(3.9), Inches(0.06), c)
    tb(sl, x + Inches(0.3), Inches(2.6), Inches(3.3), Inches(0.5), t, 22, True, c)
    tb(sl, x + Inches(0.3), Inches(3.2), Inches(3.3), Inches(1.5), d, 16, color=C_GRAY)

rrect(sl, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.5), C_CARD, C_SUBTLE)
rect(sl, Inches(0.5), Inches(5.2), Inches(0.08), Inches(1.5), C_PINK)
tb(sl, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.4), "\U0001f4a1 應用場景", 20, True, C_PINK)
tb(sl, Inches(0.9), Inches(5.8), Inches(11.5), Inches(0.7),
    "行銷素材自動化：Text Generator 撰寫社群貼文 → Image Generator 生成配圖 → Document Maker 打包成品牌簡報 → 一次性產出整月素材",
    16, color=C_GRAY)

# ════════════════════════════════════════════
# SLIDE 8: 資料處理類 Skills
# ════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "資料處理類 Skills", "分析 · 視覺化 · 資料庫操作")

skills_detail4 = [
    ("\U0001f4ca Data Analyzer", "• 數據清洗與預處理\n• 統計分析與趨勢預測\n• 異常值偵測\n• 關聯性與迴歸分析", C_PRIMARY),
    ("\U0001f4c8 Chart Generator", "• 自動生成各類型圖表\n• 長條圖/折線圖/圓餅圖\n• 互動式網頁圖表\n• 報表級排版輸出", C_ACCENT),
    ("\U0001f4be DB Query", "• 自然語言轉 SQL\n• 多資料庫支援\n• 查詢結果自動分析\n• Schema 探索與文件化", C_PINK),
]

for i, (t, d, c) in enumerate(skills_detail4):
    x = Inches(0.5 + i * 4.2)
    rrect(sl, x, Inches(2.4), Inches(3.9), Inches(2.5), C_CARD, C_SUBTLE)
    rect(sl, x, Inches(2.4), Inches(3.9), Inches(0.06), c)
    tb(sl, x + Inches(0.3), Inches(2.6), Inches(3.3), Inches(0.5), t, 22, True, c)
    tb(sl, x + Inches(0.3), Inches(3.2), Inches(3.3), Inches(1.5), d, 16, color=C_GRAY)

rrect(sl, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.5), C_CARD, C_SUBTLE)
rect(sl, Inches(0.5), Inches(5.2), Inches(0.08), Inches(1.5), C_YELLOW)
tb(sl, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.4), "\U0001f4a1 應用場景", 20, True, C_YELLOW)
tb(sl, Inches(0.9), Inches(5.8), Inches(11.5), Inches(0.7),
    "營運報表自動化：DB Query 從資料庫抓取數據 → Data Analyzer 計算 KPI 指標 → Chart Generator 產出趨勢圖表 → 整併為營運儀表板",
    16, color=C_GRAY)

# ════════════════════════════════════════════
# SLIDE 9: 工具串接類 Skills
# ════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "工具串接類 Skills", "API · Function Calling · 通訊")

skills_detail5 = [
    ("\U0001f4e1 API Caller", "• 呼叫外部 REST / GraphQL API\n• 自動處理認證與重試\n• 請求/回應格式轉換\n• Webhook 事件監聽", C_PRIMARY),
    ("\U0001f916 Function Calling", "• LLM 自主決定調用函式\n• 參數自動填入與驗證\n• 多工具並行調度\n• 錯誤處理與回退策略", C_ACCENT),
    ("\U0001f4e8 Communication", "• Email 發送與管理\n• Slack / Teams 訊息推送\n• SMS 簡訊通知\n• 社群平台自動發文", C_PINK),
]

for i, (t, d, c) in enumerate(skills_detail5):
    x = Inches(0.5 + i * 4.2)
    rrect(sl, x, Inches(2.4), Inches(3.9), Inches(2.5), C_CARD, C_SUBTLE)
    rect(sl, x, Inches(2.4), Inches(3.9), Inches(0.06), c)
    tb(sl, x + Inches(0.3), Inches(2.6), Inches(3.3), Inches(0.5), t, 22, True, c)
    tb(sl, x + Inches(0.3), Inches(3.2), Inches(3.3), Inches(1.5), d, 16, color=C_GRAY)

rrect(sl, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.5), C_CARD, C_SUBTLE)
rect(sl, Inches(0.5), Inches(5.2), Inches(0.08), Inches(1.5), C_PRIMARY)
tb(sl, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.4), "\U0001f4a1 應用場景", 20, True, C_PRIMARY)
tb(sl, Inches(0.9), Inches(5.8), Inches(11.5), Inches(0.7),
    "客戶服務自動化：Email Skill 接收客戶來信 → Function Calling 判斷意圖並調用 CRM API → API Caller 更新資料庫 → Communication Skill 回覆客戶",
    16, color=C_GRAY)

# ════════════════════════════════════════════
# SLIDE 10: 進階 Skills
# ════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "進階 Skills", "規劃 · 記憶 · 多模態理解")

advanced = [
    ("\U0001f9e0 規劃能力", "• 複雜任務分解為子步驟\n• 動態調整執行順序\n• 依賴關係自動分析\n• 失敗重試與替代方案", C_PRIMARY),
    ("\U0001f4be 記憶管理", "• 短期記憶（會話上下文）\n• 長期記憶（向量資料庫）\n• 關鍵資訊摘要壓縮\n• 跨會話知識檢索", C_ACCENT),
    ("\U0001f310 多模態理解", "• 圖片內容辨識與描述\n• 音訊轉文字與分析\n• PDF/掃描檔結構解析\n• 影片摘要與關鍵幀提取", C_PINK),
    ("\U0001f504 自我反思", "• 執行結果自動評估\n• 錯誤模式識別與學習\n• 策略調整與行為優化\n• 持續改進工作流程", C_GREEN),
]

for i, (t, d, c) in enumerate(advanced):
    x = Inches(0.6 + i * 3.15)
    rrect(sl, x, Inches(2.4), Inches(2.95), Inches(2.8), C_CARD, C_SUBTLE)
    rect(sl, x, Inches(2.4), Inches(2.95), Inches(0.06), c)
    tb(sl, x + Inches(0.3), Inches(2.6), Inches(2.35), Inches(0.5), t, 20, True, c)
    tb(sl, x + Inches(0.3), Inches(3.2), Inches(2.35), Inches(1.8), d, 15, color=C_GRAY)

rrect(sl, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.3), C_CARD, C_SUBTLE)
rect(sl, Inches(0.6), Inches(5.5), Inches(0.08), Inches(1.3), C_PRIMARY)
tb(sl, Inches(1.0), Inches(5.6), Inches(11.3), Inches(0.4), "\U0001f4a1 進階 Skills 是 AI Agent 從「工具」邁向「自主智能」的關鍵", 20, True, C_PRIMARY)
tb(sl, Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.5),
    "具備規劃 + 記憶 + 反思能力的 Agent 可自主完成跨越多步驟、多工具的複雜任務，且能從錯誤中學習與進步。", 16, color=C_GRAY)

# ════════════════════════════════════════════
# SLIDE 11: 實戰組合應用
# ════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "實戰組合應用", "如何組合 Skills 解決真實問題")

scenarios = [
    ("\U0001f4ca 競品分析報告", "Web Search\n→ 蒐集競品資訊\n\nWeb Scraper\n→ 抓取產品規格\n\nData Analyzer\n→ 比較分析\n\nDocument Maker\n→ 產出報告", C_PRIMARY),
    ("\U0001f4b0 銷售預測系統", "DB Query\n→ 拉取歷史銷售\n\nCode Interpreter\n→ 建模預測\n\nChart Generator\n→ 趨勢圖表\n\nEmail Sender\n→ 自動發送", C_ACCENT),
    ("\U0001f3af 智慧客服機器人", "File Reader\n→ 讀取 FAQ\n\nFunction Calling\n→ CRM 查詢\n\nAPI Caller\n→ 開立工單\n\nCommunication\n→ 回覆客戶", C_PINK),
]

for i, (scenario, steps, color) in enumerate(scenarios):
    x = Inches(0.5 + i * 4.2)
    rrect(sl, x, Inches(2.4), Inches(3.9), Inches(4.2), C_CARD, C_SUBTLE)
    rect(sl, x, Inches(2.4), Inches(3.9), Inches(0.06), color)
    tb(sl, x + Inches(0.3), Inches(2.6), Inches(3.3), Inches(0.5), scenario, 22, True, color)
    tb(sl, x + Inches(0.3), Inches(3.2), Inches(3.3), Inches(3.0), steps, 15, color=C_GRAY)

# ════════════════════════════════════════════
# SLIDE 12: 結語
# ════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, C_DARK)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.08), C_PRIMARY)
rect(sl, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), C_PRIMARY)

# decorative circles
for i, (x, r, c) in enumerate([
    (Inches(9.5), Inches(2.5), C_PRIMARY),
    (Inches(10.8), Inches(1.8), C_ACCENT),
]):
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(1.0 + i*0.5), r, r)
    circ.fill.solid(); circ.fill.fore_color.rgb = c
    circ.fill.fore_color.brightness = 0.7; circ.line.fill.background()

tb(sl, Inches(1.0), Inches(1.0), Inches(11), Inches(0.5), "結語與下一步", 16, color=C_SECONDARY)
tb(sl, Inches(1.0), Inches(1.6), Inches(7.5), Inches(1.2),
    "掌握 Skills 組合\n釋放 AI Agent 真正潛力", 44, True, C_WHITE)
rect(sl, Inches(1.0), Inches(3.0), Inches(2.5), Inches(0.06), C_PRIMARY)

rrect(sl, Inches(1.0), Inches(3.4), Inches(11.3), Inches(1.8), C_CARD, C_SUBTLE)
tb(sl, Inches(1.4), Inches(3.6), Inches(10.5), Inches(0.4), "\U0001f4a1 給開發者的建議", 20, True, C_PRIMARY)
tips = [
    "\U000025b6 從「單一 Skill」開始熟悉，再慢慢組合多個 Skills",
    "\U000025b6 每個任務先思考「需要哪些資訊取得 → 如何處理 → 如何輸出」",
    "\U000025b6 善用 Planning Skill 讓 Agent 自主編排工作流程",
    "\U000025b6 持續迭代 Prompt，讓 Agent 更精準選擇與使用 Skills",
]
for i, tip in enumerate(tips):
    tb(sl, Inches(1.6), Inches(4.1 + i * 0.35), Inches(10.3), Inches(0.35), tip, 15, color=C_GRAY)

# skill icons row
icons = ["\U0001f50d", "\U0001f4bb", "\U0001f3a8", "\U0001f4ca", "\U0001f4e1", "\U0001f9e0", "\U0001f4be", "\U0001f310"]
for i, icon in enumerate(icons):
    rrect(sl, Inches(1.0 + i * 1.5), Inches(5.5), Inches(1.2), Inches(0.7), C_SUBTLE)
    tb(sl, Inches(1.0 + i * 1.5), Inches(5.52), Inches(1.2), Inches(0.66), icon, 28, align=PP_ALIGN.CENTER)

tb(sl, Inches(1.0), Inches(6.5), Inches(5), Inches(0.4), "聯絡我們：skills@aiagent-ecosystem.com", 14, color=C_SECONDARY)

# ── Save ──
output = "AI_Agent_Skills_介紹.pptx"
prs.save(output)
print(f"✅ 簡報已生成: {output}")
